#!/usr/bin/env python3
"""Corporate PPT renderer for corporate-ppt-bulk (HTML-first pipeline).

Usage: python3 render.py <plan.json> <output.pptx>

plan.json shape (same as the corporate-ppt Claude Skill's, notes inline
per slide rather than a separate script-generation step):
{
  "unit": "...", "epigrafe": "...", "afo": "...",
  "slides": [
    {"n": 1, "section": "titulo", "fields": {"title": "..."}, "notes": "..."},
    {"n": 2, "section": "inicio", "fields": {"icon": "lightbulb", "promise": "..."}, "notes": "..."},
    {"n": 3, "section": "concepto", "variant": "numero_hero", "fields": {...}, "notes": "..."},
    {"n": 4, "section": "puntos_clave", "variant": "flujo_pasos", "fields": {...}, "notes": "..."},
    {"n": 5, "section": "resumen", "fields": {"title": "...", "items": [...]}, "notes": "..."},
    {"n": 6, "section": "cierre", "fields": {"title": "Thank you"}, "notes": "..."}
  ]
}

Each slide renders as real HTML/CSS (icons.py + templates.py + slides.css —
same design system as the .claude/skills/corporate-ppt Claude Skill),
screenshotted via headless Chromium, then reassembled with titles/captions
as REAL editable pptx text boxes overlaid on top — only decorative art
(the giant hero number, card shadows, gradients) is baked into the
background image. Rubik/Lato are embedded into the .pptx afterward so
those text boxes render correctly even without the fonts installed.

Called directly from jobs.py (assemble_pptx/embed_fonts imported, not
shelled out to) — the __main__ CLI entry point below is kept for
standalone testing.
"""
import json
import os
import re
import sys
import tempfile
import zipfile

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from templates import render_slide  # noqa: E402

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
FONTS_DIR = os.path.join(HERE, "..", "fonts")
SLIDE_W_IN, SLIDE_H_IN = 10.0, 5.625
PX_W, PX_H = 1920, 1080
PX_PER_IN = PX_W / SLIDE_W_IN      # 192
PT_PER_PX = 72.0 / PX_PER_IN       # 0.375
SCALE = 2                          # device_scale_factor for screenshots

FONTS = [
    {"typeface": "Rubik", "style": "regular", "file": "Rubik-Regular.ttf"},
    {"typeface": "Rubik", "style": "bold", "file": "Rubik-Bold.ttf"},
    {"typeface": "Lato", "style": "regular", "file": "Lato-Regular.ttf"},
    {"typeface": "Lato", "style": "bold", "file": "Lato-Bold.ttf"},
]

EXTRACT_JS = """
() => {
  const els = document.querySelectorAll('.pptx-text');
  const results = [];
  els.forEach((el) => {
    if (!el.textContent.trim()) { el.style.visibility = 'hidden'; return; }
    const rect = el.getBoundingClientRect();
    const cs = getComputedStyle(el);
    const runs = [];
    el.childNodes.forEach((node) => {
      if (node.nodeType === Node.TEXT_NODE) {
        if (!node.textContent.trim()) return;
        runs.push({ text: node.textContent, color: cs.color, weight: cs.fontWeight, family: cs.fontFamily, size: cs.fontSize });
      } else if (node.nodeType === Node.ELEMENT_NODE) {
        if (!node.textContent.trim()) return;
        const ccs = getComputedStyle(node);
        runs.push({ text: node.textContent, color: ccs.color, weight: ccs.fontWeight, family: ccs.fontFamily, size: ccs.fontSize });
      }
    });
    if (runs.length) {
      results.push({ left: rect.left, top: rect.top, width: rect.width, height: rect.height, align: cs.textAlign, runs });
    }
    el.style.visibility = 'hidden';
  });
  return results;
}
"""


def find_chrome():
    candidates = [
        os.environ.get("PLAYWRIGHT_CHROMIUM_PATH"),
        # System-installed Chromium, if present — otherwise falls through
        # to Playwright's own bundled download (see PLAYWRIGHT_BROWSERS_PATH
        # below), which is the normal case on the Debian-based production image.
        "/usr/bin/chromium-browser",
        "/usr/bin/chromium",
    ]
    pw_dir = os.environ.get("PLAYWRIGHT_BROWSERS_PATH", "/opt/pw-browsers")
    if os.path.isdir(pw_dir):
        for name in sorted(os.listdir(pw_dir)):
            if name.startswith("chromium-"):
                candidates.append(os.path.join(pw_dir, name, "chrome-linux", "chrome"))
    for c in candidates:
        if c and os.path.isfile(c):
            return c
    return None  # let Playwright fall back to its own resolution


def build_html(slide):
    cls, style, inner = render_slide(slide)
    fonts_dir = os.path.abspath(FONTS_DIR)
    slides_css = open(os.path.join(HERE, "slides.css"), encoding="utf-8").read()
    return f"""<!doctype html><html><head><meta charset="utf-8"><style>
  @font-face {{ font-family:"Rubik"; src:url("file://{fonts_dir}/Rubik-Regular.ttf") format("truetype"); font-weight:400; }}
  @font-face {{ font-family:"Rubik"; src:url("file://{fonts_dir}/Rubik-Bold.ttf") format("truetype"); font-weight:700; }}
  @font-face {{ font-family:"Lato"; src:url("file://{fonts_dir}/Lato-Regular.ttf") format("truetype"); font-weight:400; }}
  @font-face {{ font-family:"Lato"; src:url("file://{fonts_dir}/Lato-Bold.ttf") format("truetype"); font-weight:700; }}
  {slides_css}
</style></head><body>
<div class="slide-canvas {cls}"{style}>{inner}</div>
</body></html>"""


def rgb_to_hex(rgb_str):
    m = re.match(r"rgba?\((\d+),\s*(\d+),\s*(\d+)", rgb_str)
    if not m:
        return "202020"
    return "".join(f"{int(x):02X}" for x in m.groups())


def first_family(css_family):
    return css_family.split(",")[0].strip().strip('"').strip("'")


def capture_slide(page, slide, tmp_dir, index):
    html_path = os.path.join(tmp_dir, f"slide-{index}.html")
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(build_html(slide))
    page.goto(f"file://{html_path}")
    page.wait_for_timeout(60)  # let @font-face finish applying
    text_boxes = page.evaluate(EXTRACT_JS)
    img_path = os.path.join(tmp_dir, f"slide-{index}.png")
    page.screenshot(path=img_path)
    return img_path, text_boxes


def assemble_pptx(plan, tmp_dir, out_path):
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W_IN * 914400))
    prs.slide_height = Emu(int(SLIDE_H_IN * 914400))
    prs.core_properties.title = plan.get("epigrafe") or plan.get("unit") or ""

    chrome = find_chrome()
    with sync_playwright() as p:
        launch_kwargs = {"args": ["--no-sandbox"]}
        if chrome:
            launch_kwargs["executable_path"] = chrome
        browser = p.chromium.launch(**launch_kwargs)
        page = browser.new_page(viewport={"width": PX_W, "height": PX_H}, device_scale_factor=SCALE)

        for i, slide_data in enumerate(plan["slides"]):
            img_path, text_boxes = capture_slide(page, slide_data, tmp_dir, i)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))

            for box in text_boxes:
                tb = slide.shapes.add_textbox(
                    Inches(box["left"] / PX_PER_IN), Inches(box["top"] / PX_PER_IN),
                    Inches(box["width"] / PX_PER_IN), Inches(box["height"] / PX_PER_IN),
                )
                tf = tb.text_frame
                tf.word_wrap = True
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                para = tf.paragraphs[0]
                para.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(box["align"], PP_ALIGN.LEFT)
                for r in box["runs"]:
                    run = para.add_run()
                    run.text = r["text"]
                    run.font.name = first_family(r["family"])
                    run.font.size = Pt(round(float(r["size"].replace("px", "")) * PT_PER_PX))
                    run.font.bold = int(r["weight"]) >= 700
                    run.font.color.rgb = RGBColor.from_string(rgb_to_hex(r["color"]))

            notes = slide_data.get("notes")
            if notes and notes.strip():
                slide.notes_slide.notes_text_frame.text = notes

        browser.close()

    prs.save(out_path)


def embed_fonts(pptx_path):
    """Embed Rubik/Lato into the .pptx so editable text boxes render
    correctly even without those fonts installed locally (direct
    zipfile/OOXML manipulation, no external tool)."""
    font_bufs = [{**f, "buffer": open(os.path.join(FONTS_DIR, f["file"]), "rb").read()} for f in FONTS]

    with zipfile.ZipFile(pptx_path, "r") as zin:
        names = zin.namelist()
        content_types = zin.read("[Content_Types].xml").decode("utf-8")
        presentation = zin.read("ppt/presentation.xml").decode("utf-8")
        rels = zin.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
        others = {n: zin.read(n) for n in names if n not in
                  ("[Content_Types].xml", "ppt/presentation.xml", "ppt/_rels/presentation.xml.rels")}

    if 'Extension="fntdata"' not in content_types:
        content_types = content_types.replace(
            "</Types>", '<Default Extension="fntdata" ContentType="application/x-fontdata"/></Types>')

    existing_ids = [int(m) for m in re.findall(r'Id="rId(\d+)"', rels)]
    next_id = (max(existing_ids) + 1) if existing_ids else 1

    font_parts = {}
    with_rel_ids = []
    for i, font in enumerate(font_bufs):
        rid = f"rId{next_id + i}"
        part_name = f"font{next_id + i}.fntdata"
        font_parts[f"ppt/fonts/{part_name}"] = font["buffer"]
        rels = rels.replace(
            "</Relationships>",
            f'<Relationship Id="{rid}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/font" Target="fonts/{part_name}"/></Relationships>')
        with_rel_ids.append({**font, "rId": rid})

    if "embedTrueTypeFonts=" not in presentation:
        presentation = presentation.replace("<p:presentation ", '<p:presentation embedTrueTypeFonts="1" ', 1)

    by_typeface = {}
    for f in with_rel_ids:
        by_typeface.setdefault(f["typeface"], {})[f["style"]] = f["rId"]
    embedded_font_lst = "<p:embeddedFontLst>" + "".join(
        f'<p:embeddedFont><p:font typeface="{typeface}"/>' +
        "".join(f'<p:{style} r:id="{rid}"/>' for style, rid in styles.items()) +
        "</p:embeddedFont>"
        for typeface, styles in by_typeface.items()
    ) + "</p:embeddedFontLst>"

    if re.search(r"<p:notesSz\b[^/]*/>", presentation):
        presentation = re.sub(r"(<p:notesSz\b[^/]*/>)", r"\1" + embedded_font_lst, presentation, count=1)
    else:
        presentation = presentation.replace("</p:presentation>", embedded_font_lst + "</p:presentation>")

    with zipfile.ZipFile(pptx_path, "w", zipfile.ZIP_DEFLATED) as zout:
        zout.writestr("[Content_Types].xml", content_types)
        zout.writestr("ppt/presentation.xml", presentation)
        zout.writestr("ppt/_rels/presentation.xml.rels", rels)
        for name, data in others.items():
            zout.writestr(name, data)
        for name, data in font_parts.items():
            zout.writestr(name, data)


def main():
    input_path, output_path = sys.argv[1], sys.argv[2]
    with open(input_path, "r", encoding="utf-8") as f:
        plan = json.load(f)

    with tempfile.TemporaryDirectory() as tmp_dir:
        assemble_pptx(plan, tmp_dir, output_path)

    try:
        embed_fonts(output_path)
    except Exception as e:  # non-fatal: ship without embedded fonts rather than fail
        print(f"WARNING: font embedding failed ({e}); shipping without embedded fonts.", file=sys.stderr)


if __name__ == "__main__":
    main()
