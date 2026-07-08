#!/usr/bin/env python3
"""
Corporate PPT renderer — EDUCA EDTECH Group design system.

Replaces the pptxgenjs-based renderCorporateSlide() for the PPT+Script
flow. Reason: pptxgenjs has no native multi-stop gradient fill (open
upstream issue since 2017), so the brand gradient (§6.2/§6.3 of the
manual) could only be approximated with solid segments. python-pptx
exposes the underlying DrawingML <a:gsLst> directly, so we can write the
real 5-stop brand gradient.

Usage: python3 render_ppt.py <input.json> <output.pptx>

input.json shape: { "plan": {...slide plan...}, "scripts": [...generated
script/productionNotes per slide, keyed by "n"...] }

Font embedding is NOT done here — server/index.js post-processes the
resulting .pptx (embedFonts()) regardless of which renderer produced it.
"""
import json
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

EMU_PER_IN = 914400
SLIDE_W = 10.0
SLIDE_H = 5.625

# Icon set: Material Symbols (Outlined), Apache 2.0 — pre-rasterized to
# PNG in brand colors (see server/icons/LICENSE.txt). GPT-4o picks from
# this exact set; anything else is silently skipped (no icon rendered).
ICONS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "icons")
ICON_COLORS = {"FFFFFF", "244A80", "2E7ABE", "963058", "E96A73", "60BFB8", "202020"}

# Brand gradient stops (teal -> blue-m -> blue-d -> burdeos -> rosa)
GRADIENT_STOPS = [
    (0, "60BFB8"),
    (25000, "2E7ABE"),
    (50000, "244A80"),
    (80000, "963058"),
    (100000, "E96A73"),
]

SECTION_CONFIG = {
    "title":        {"bg": "60BFB8", "accent": "60BFB8", "dark": True,  "label": None,
                      "gradient": True},
    "entrada":      {"bg": "FFFFFF", "accent": "60BFB8", "dark": False, "label": "Introduction"},
    "conceptos":    {"bg": "FFFFFF", "accent": "244A80", "dark": False, "label": "Key Concepts"},
    "puntos_clave": {"bg": "FFFFFF", "accent": "2E7ABE", "dark": False, "label": "Key Points"},
    "resumen":      {"bg": "963058", "accent": "963058", "dark": True,  "label": None},
    "cierre":       {"bg": "E96A73", "accent": "E96A73", "dark": True,  "label": None,
                      "gradient": True},
}
# Microlearning format uses "inicio"/"concepto" (singular, one slide per
# epigraph) instead of "entrada"/"conceptos" — same visual treatment.
SECTION_CONFIG["inicio"] = SECTION_CONFIG["entrada"]
SECTION_CONFIG["concepto"] = SECTION_CONFIG["conceptos"]


def rgb(hex_color):
    return RGBColor.from_string(hex_color)


def set_gradient_fill(fill, stops=GRADIENT_STOPS, angle_deg=0):
    """Write a real multi-stop DrawingML gradient (python-pptx's public
    API only supports the default 2 stops — this manipulates the
    underlying <a:gsLst> XML directly to add the rest)."""
    fill.gradient()
    grad_fill_el = fill._xPr.find(qn("a:gradFill"))
    gs_lst = grad_fill_el.find(qn("a:gsLst"))
    for child in list(gs_lst):
        gs_lst.remove(child)
    for pos, hex_color in stops:
        gs = gs_lst.makeelement(qn("a:gs"), {"pos": str(pos)})
        clr = gs.makeelement(qn("a:srgbClr"), {"val": hex_color})
        gs.append(clr)
        gs_lst.append(gs)
    lin = grad_fill_el.find(qn("a:lin"))
    if lin is not None:
        lin.set("ang", str(angle_deg * 60000))


def add_rect(slide, x, y, w, h, color, shape_type=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(color)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_gradient_bar(slide, y, h):
    shp = add_rect(slide, 0, y, SLIDE_W, h, "FFFFFF")
    set_gradient_fill(shp.fill, angle_deg=0)


def add_text(slide, x, y, w, h, text, font="Lato", size=14, bold=False, italic=False,
             color="202020", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, char_spacing=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return box


def add_icon(slide, icon_name, x, y, size, color="202020"):
    """Place a pre-baked icon PNG. No-op (skips silently) if the
    (icon_name, color) pair isn't one of the pre-rasterized assets —
    GPT-4o is instructed to only use the documented icon set, but a
    bad/hallucinated name shouldn't break slide generation."""
    if not icon_name or color.upper() not in ICON_COLORS:
        return
    path = os.path.join(ICONS_DIR, f"{icon_name}__{color.upper()}.png")
    if not os.path.isfile(path):
        return
    slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(size), Inches(size))


def _as_dict(item):
    """GPT-4o sometimes sends a plain string where the schema documents
    an object (e.g. ["Discover", "Apply"] instead of [{"title": ...}]).
    Coerce so the renderer still shows something instead of a blank
    panel — rather than silently dropping a slide's only visual."""
    if isinstance(item, dict):
        return item
    if isinstance(item, str):
        return {"title": item, "text": "", "label": item, "step": item, "word": item}
    return {}


def _text_of(item, *keys):
    d = _as_dict(item)
    for k in keys:
        v = d.get(k)
        if v:
            return str(v)
    return ""


# ── Graphic renderers (mirrors server/index.js renderGraphic()) ────────
def render_graphic(slide, gtype, data, box):
    x, y, w, h = box["x"], box["y"], box["w"], box["h"]
    if not isinstance(data, dict):
        data = {}

    rendered_count = 0

    if gtype == "text_only":
        text = data.get("text", "")
        add_text(slide, x, y, w, h, text, font="Lato", size=13,
                  italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rendered_count += 1 if text else 0

    elif gtype == "before_after":
        before_text = data.get("beforeText", "")
        after_text = data.get("afterText", "")
        add_text(slide, x, y, w, 1,
                  f"{data.get('beforeLabel', 'Before')}: {before_text}",
                  size=12, color="963058")
        add_text(slide, x, y + 1.2, w, 1.2,
                  f"{data.get('afterLabel', 'After')}: {after_text}",
                  size=12, color="202020")
        rendered_count += 1 if (before_text or after_text) else 0

    elif gtype == "smart_grid":
        items = [_as_dict(i) for i in data.get("items") or []]
        highlight = data.get("highlightLetter")
        for i, item in enumerate(items):
            col, row = i % 2, i // 2
            cx = x + col * (w / 2)
            cy = y + row * 0.9
            letter = item.get("letter") or (_text_of(item, "title", "word", "text")[:1])
            word = item.get("word") or _text_of(item, "title", "text")
            is_hl = letter == highlight
            add_rect(slide, cx, cy, w / 2 - 0.05, 0.85, "963058" if is_hl else "F0F0F0")
            add_text(slide, cx, cy, w / 2 - 0.05, 0.55, letter,
                      font="Rubik", size=22, bold=True,
                      color="FFFFFF" if is_hl else "666666",
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, cx, cy + 0.55, w / 2 - 0.05, 0.3, word,
                      size=8, color="666666", align=PP_ALIGN.CENTER)
            rendered_count += 1

    elif gtype == "data_table":
        columns = data.get("columns") or []
        raw_rows = data.get("rows") or []
        # Coerce rows given as {col: val} dicts instead of positional arrays.
        rows = [
            [r.get(c, "") for c in columns] if isinstance(r, dict) else r
            for r in raw_rows
        ]
        n_rows = len(rows) + 1
        n_cols = max(len(columns), 1)
        gfx = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(min(h, 0.4 * n_rows)))
        table = gfx.table
        for c, col_name in enumerate(columns):
            cell = table.cell(0, c)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb("202020")
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.bold = True
            run.font.color.rgb = rgb("FFFFFF")
            run.font.size = Pt(11)
        for r, row_data in enumerate(rows):
            for c, cell_val in enumerate(row_data):
                cell = table.cell(r + 1, c)
                cell.text = str(cell_val)
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
        rendered_count += len(rows)

    elif gtype == "three_node_sequence":
        nodes = data.get("nodes") or []
        n = max(len(nodes), 1)
        slot_w = w / n
        node_w = slot_w - 0.18
        for i, node in enumerate(nodes):
            label = node if isinstance(node, str) else _text_of(node, "label", "title", "text", "step")
            nx = x + i * slot_w
            ny = y + h / 2 - 0.35
            add_rect(slide, nx, ny, node_w, 0.7, "244A80", MSO_SHAPE.ROUNDED_RECTANGLE)
            add_text(slide, nx, ny, node_w, 0.7, label, size=11, color="FFFFFF",
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if i < n - 1:
                add_text(slide, nx + node_w, y + h / 2 - 0.25, slot_w - node_w, 0.5, "→",
                          size=16, color="963058", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            rendered_count += 1

    elif gtype == "numbered_list":
        items = data.get("items") or []
        row_h = h / max(len(items), 1)
        for i, item in enumerate(items):
            text = item if isinstance(item, str) else _text_of(item, "text", "title", "label", "step")
            iy = y + i * row_h
            add_rect(slide, x, iy + row_h / 2 - 0.18, 0.36, 0.36, "2E7ABE", MSO_SHAPE.OVAL)
            add_text(slide, x, iy + row_h / 2 - 0.18, 0.36, 0.36, str(i + 1), font="Rubik",
                      size=13, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, x + 0.5, iy, w - 0.5, row_h, text, size=12, anchor=MSO_ANCHOR.MIDDLE)
            rendered_count += 1

    elif gtype == "validation_flow":
        steps = data.get("steps") or []
        row_h = h / max(len(steps), 1)
        for i, step in enumerate(steps):
            text = step if isinstance(step, str) else _text_of(step, "text", "title", "label", "step")
            iy = y + i * row_h
            add_text(slide, x, iy, 0.4, row_h, "✓", size=16, bold=True, color="60BFB8",
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, x + 0.45, iy, w - 0.45, row_h, text, size=12, anchor=MSO_ANCHOR.MIDDLE)
            rendered_count += 1

    elif gtype == "pillar_columns":
        cols = [_as_dict(c) for c in data.get("columns") or []]
        n = max(len(cols), 1)
        slot_w = w / n
        col_w = slot_w - 0.1
        for i, col in enumerate(cols):
            cx = x + i * slot_w
            add_rect(slide, cx, y, col_w, h, "F0F0F0")
            icon = col.get("icon")
            title_y = y + 0.1
            if icon:
                icon_size = 0.4
                add_icon(slide, icon, cx + col_w / 2 - icon_size / 2, y + 0.12, icon_size, "963058")
                title_y = y + 0.55
            add_text(slide, cx + 0.08, title_y, col_w - 0.16, 0.4, _text_of(col, "title", "label"),
                      font="Rubik", size=11, bold=True, color="963058", align=PP_ALIGN.CENTER)
            add_text(slide, cx + 0.08, title_y + 0.45, col_w - 0.16, h - (title_y - y) - 0.55,
                      _text_of(col, "text", "description"), size=10, color="202020", align=PP_ALIGN.CENTER)
            rendered_count += 1

    elif gtype == "icon_grid":
        # 2-column grid of icon + title + short description cards — the
        # dominant graphic pattern for concept/key-point slides (icon
        # cards), heavier than pillar_columns which is a single row.
        items = [_as_dict(i) for i in data.get("items") or []]
        n = max(len(items), 1)
        cols = 2 if n > 1 else 1
        rows = (n + cols - 1) // cols
        cell_w = w / cols
        cell_h = h / rows
        for i, item in enumerate(items):
            col, row = i % cols, i // cols
            cx = x + col * cell_w
            cy = y + row * cell_h
            pad = 0.06
            add_rect(slide, cx + pad, cy + pad, cell_w - 2 * pad, cell_h - 2 * pad, "F5F5F5", MSO_SHAPE.ROUNDED_RECTANGLE)
            icon = item.get("icon")
            text_y = cy + pad + 0.08
            if icon:
                icon_size = 0.34
                add_icon(slide, icon, cx + pad + 0.12, cy + pad + 0.08, icon_size, "244A80")
                text_y = cy + pad + 0.08
                title_x = cx + pad + 0.12 + icon_size + 0.1
            else:
                title_x = cx + pad + 0.12
            title_w = cell_w - 2 * pad - (title_x - (cx + pad))
            add_text(slide, title_x, text_y, title_w, 0.3, _text_of(item, "title", "label"),
                      font="Rubik", size=10, bold=True, color="244A80", anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, cx + pad + 0.12, cy + pad + 0.5, cell_w - 2 * pad - 0.24, cell_h - 2 * pad - 0.6,
                      _text_of(item, "text", "description"), size=9, color="202020")
            rendered_count += 1

    else:
        # "none" and any unrecognized type: no-op
        return

    if rendered_count == 0:
        # The panel/container was drawn (has_graphic was true) but the
        # graphicData didn't contain anything render_graphic recognized —
        # almost always means GPT-4o's JSON shape didn't match the
        # documented schema for this graphicType. Surface it so it shows
        # up in server logs instead of just an empty panel in the deck.
        print(f"WARNING: graphicType '{gtype}' rendered 0 items — graphicData was: {json.dumps(data)[:300]}",
              file=sys.stderr)


# ── Corporate slide renderer (mirrors renderCorporateSlide) ────────────
def render_slide(prs, slide_data):
    section = slide_data.get("section")
    cfg = SECTION_CONFIG.get(section, SECTION_CONFIG["conceptos"])
    bullets = (slide_data.get("bullets") or [])[:3]
    title = slide_data.get("title") or ""
    subtitle = slide_data.get("subtitle")

    graphic_type = slide_data.get("graphicType")
    if graphic_type == "none":
        graphic_type = None
    legacy_graphic = None
    if not graphic_type and slide_data.get("graphic") and slide_data["graphic"] != "none":
        legacy_graphic = slide_data["graphic"]
    has_graphic = bool(graphic_type or legacy_graphic)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    if cfg["dark"]:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(cfg["bg"])
        if cfg.get("gradient"):
            # Cierre bookends title with the gradient reversed (rosa→teal
            # instead of teal→rosa) — same brand ramp, opposite direction.
            positions = [pos for pos, _ in GRADIENT_STOPS]
            colors = [color for _, color in GRADIENT_STOPS]
            if section == "cierre":
                colors = list(reversed(colors))
            set_gradient_fill(slide.background.fill, stops=list(zip(positions, colors)), angle_deg=45)

        add_rect(slide, 0, 0, SLIDE_W, 0.06, "FFFFFF")

        add_text(slide, 0.8, 1.4, 8.4, 1.8, title, font="Rubik", size=38, color="FFFFFF",
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        if subtitle:
            add_text(slide, 0.8, 3.2, 8.4, 0.7, subtitle, font="Lato", size=18, color="FFFFFF",
                      align=PP_ALIGN.CENTER)

        for i, bullet in enumerate(bullets):
            add_text(slide, 1.5, 2.8 + i * 0.72, 7, 0.65, f"• {bullet}",
                      font="Lato", size=18, color="FFFFFF", align=PP_ALIGN.CENTER)

        add_gradient_bar(slide, 5.43, 0.2)

    else:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(cfg["bg"])

        add_rect(slide, 0, 0, 0.14, SLIDE_H, cfg["accent"])

        if cfg["label"]:
            add_text(slide, 0.36, 0.2, 9.3, 0.32, cfg["label"].upper(), font="Lato", size=10,
                      bold=True, color=cfg["accent"])

        add_text(slide, 0.36, 0.52, 9.3, 1.0, title, font="Rubik", size=28, color=cfg["accent"])

        add_rect(slide, 0.36, 1.52, 9.28, 0.018, "E0E0E0")

        bullet_w = 5.7 if has_graphic else 9.06
        bullet_y = 1.7
        spacing = 1.1 if len(bullets) <= 2 else 0.9
        for i, bullet in enumerate(bullets):
            add_rect(slide, 0.36, bullet_y + i * spacing + 0.27, 0.07, 0.07, cfg["accent"])
            add_text(slide, 0.58, bullet_y + i * spacing, bullet_w, spacing - 0.05, bullet,
                      font="Lato", size=18, color="202020", anchor=MSO_ANCHOR.MIDDLE)

        if has_graphic:
            panel_x, panel_y, panel_w, panel_h = 6.3, 1.7, 3.3, 3.4
            add_rect(slide, panel_x, panel_y, panel_w, panel_h, "F5F5F5", MSO_SHAPE.ROUNDED_RECTANGLE)
            inner = {"x": panel_x + 0.25, "y": panel_y + 0.25, "w": panel_w - 0.5, "h": panel_h - 0.5}
            if graphic_type:
                render_graphic(slide, graphic_type, slide_data.get("graphicData") or {}, inner)
            else:
                add_text(slide, inner["x"], inner["y"], inner["w"], inner["h"], legacy_graphic,
                          font="Lato", size=12, color=cfg["accent"], italic=True,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        add_gradient_bar(slide, 5.43, 0.2)

    return slide


def main():
    input_path, output_path = sys.argv[1], sys.argv[2]
    with open(input_path, "r", encoding="utf-8") as f:
        payload = json.load(f)

    plan = payload["plan"]
    scripts_by_n = {s["n"]: s for s in payload.get("scripts") or []}

    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W * EMU_PER_IN))
    prs.slide_height = Emu(int(SLIDE_H * EMU_PER_IN))
    prs.core_properties.title = plan.get("unit", "")

    for slide_data in plan.get("slides", []):
        slide = render_slide(prs, slide_data)
        script_data = scripts_by_n.get(slide_data.get("n"), {})
        notes = "\n\n---\n".join(
            [v for v in [script_data.get("script"), script_data.get("productionNotes")] if v]
        )
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(output_path)


if __name__ == "__main__":
    main()
